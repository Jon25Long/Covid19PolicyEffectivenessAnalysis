"""
Automated PowerPoint Presentation Generator for Jupyter Notebooks

This module automatically extracts content, figures, and results from
Jupyter notebooks and generates comprehensive PowerPoint presentations.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import re
from PIL import Image
import io


class NotebookPresenter:
    """
    Automatically generates PowerPoint presentations from Jupyter notebooks.
    
    Usage:
        presenter = NotebookPresenter()
        presenter.auto_extract_from_notebook()
        output_path = presenter.generate_presentation()
    """
    
    def __init__(self, notebook_path=None, output_dir=None):
        """
        Initialize the presenter.
        
        Args:
            notebook_path: Path to the notebook file (auto-detected if None)
            output_dir: Output directory for presentation (default: output/presentations/)
        """
        # Auto-detect notebook if not provided
        if notebook_path is None:
            notebook_path = self._detect_current_notebook()
        
        self.notebook_path = Path(notebook_path)
        self.notebook_filename = self.notebook_path.stem
        
        # Set up output directory
        project_root = self.notebook_path.parent.parent
        if output_dir is None:
            output_dir = project_root / 'output' / 'presentations'
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Initialize storage
        self.extracted_content = {}
        self.extracted_figures = []
        self.statistical_results = {}
        
    def _detect_current_notebook(self):
        """Auto-detect the current notebook from the calling context."""
        import inspect
        import os
        
        # Try to get the current notebook from IPython
        try:
            from IPython import get_ipython
            ipython = get_ipython()
            if ipython is not None:
                # Try to get notebook name from IPython history
                notebook_name = ipython.user_ns.get('__vsc_ipynb_file__', None)
                if notebook_name:
                    return notebook_name
        except:
            pass
        
        # Fallback: look for .ipynb files in parent directory
        frame = inspect.currentframe()
        if frame and frame.f_back and frame.f_back.f_back:
            caller_file = frame.f_back.f_back.f_globals.get('__file__', None)
            if caller_file:
                caller_dir = Path(caller_file).parent
                notebooks = list(caller_dir.glob('*.ipynb'))
                if notebooks:
                    return str(notebooks[0])
        
        # Last resort: current directory
        notebooks = list(Path('.').glob('*.ipynb'))
        if notebooks:
            return str(notebooks[0])
        
        raise FileNotFoundError("Could not auto-detect notebook. Please provide notebook_path parameter.")
    
    def auto_extract_from_notebook(self):
        """
        Automatically extract all content from the notebook.
        This includes text, figures, and statistical results.
        """
        with open(self.notebook_path, 'r', encoding='utf-8') as f:
            notebook = json.load(f)
        
        # Extract content from cells
        for i, cell in enumerate(notebook['cells']):
            cell_type = cell['cell_type']
            
            if cell_type == 'markdown':
                self._extract_markdown_content(cell, i)
            elif cell_type == 'code':
                self._extract_code_outputs(cell, i)
        
        print(f"✓ Extracted content from {self.notebook_filename}")
        print(f"  - {len(self.extracted_content)} text sections")
        print(f"  - {len(self.extracted_figures)} figures")
        print(f"  - {len(self.statistical_results)} statistical results")
    
    def _extract_markdown_content(self, cell, index):
        """Extract text content from markdown cells."""
        source = ''.join(cell['source'])
        
        # Look for headers to organize content
        if source.strip().startswith('#'):
            # Extract section title
            lines = source.split('\n')
            title_line = lines[0]
            title = title_line.lstrip('#').strip()
            content = '\n'.join(lines[1:]).strip()
            
            self.extracted_content[f"section_{index}"] = {
                'title': title,
                'content': content,
                'type': 'section'
            }
    
    def _extract_code_outputs(self, cell, index):
        """Extract figures and results from code cell outputs."""
        if 'outputs' not in cell:
            return
        
        for output in cell['outputs']:
            # Extract images
            if 'data' in output:
                if 'image/png' in output['data']:
                    self.extracted_figures.append({
                        'index': index,
                        'type': 'png',
                        'data': output['data']['image/png']
                    })
                elif 'image/jpeg' in output['data']:
                    self.extracted_figures.append({
                        'index': index,
                        'type': 'jpeg',
                        'data': output['data']['image/jpeg']
                    })
            
            # Extract text output that looks like statistical results
            if output.get('output_type') == 'stream' and output.get('name') == 'stdout':
                text = ''.join(output['text'])
                if any(keyword in text.lower() for keyword in ['correlation', 'p-value', 'r-squared', 'regression']):
                    self.statistical_results[f"stats_{index}"] = text
    
    def generate_presentation(self, title=None):
        """
        Generate the PowerPoint presentation.
        
        Args:
            title: Presentation title (auto-generated if None)
            
        Returns:
            Path to the generated presentation file
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Generate title if not provided
        if title is None:
            title = self.notebook_filename.replace('_', ' ').replace('-', ' ').title()
        
        # Title slide
        self._add_title_slide(prs, title)
        
        # Add content slides
        for key, content in self.extracted_content.items():
            if content['type'] == 'section' and content['content']:
                self._add_content_slide(prs, content['title'], content['content'])
        
        # Add figure slides
        for fig in self.extracted_figures:
            self._add_figure_slide(prs, fig)
        
        # Add statistical results slides
        for key, stats in self.statistical_results.items():
            self._add_stats_slide(prs, stats)
        
        # Save presentation
        output_path = self.output_dir / f"{self.notebook_filename}_presentation.pptx"
        prs.save(str(output_path))
        
        return output_path
    
    def _add_title_slide(self, prs, title):
        """Add title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"Automated Analysis Report\nGenerated from: {self.notebook_filename}"
    
    def _add_content_slide(self, prs, title, content):
        """Add a content slide with title and text."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        text_frame = body_shape.text_frame
        text_frame.text = content[:1000]  # Limit text length
    
    def _add_figure_slide(self, prs, fig_data):
        """Add a slide with a figure."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Decode base64 image data
        import base64
        img_data = base64.b64decode(fig_data['data'])
        
        # Save temporarily and add to slide
        temp_path = self.output_dir / f"temp_fig_{fig_data['index']}.png"
        with open(temp_path, 'wb') as f:
            f.write(img_data)
        
        # Add image to slide
        left = Inches(1)
        top = Inches(1)
        pic = slide.shapes.add_picture(str(temp_path), left, top, width=Inches(8))
        
        # Clean up temp file
        temp_path.unlink()
    
    def _add_stats_slide(self, prs, stats_text):
        """Add a slide with statistical results."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = "Statistical Results"
        text_frame = body_shape.text_frame
        text_frame.text = stats_text[:1000]  # Limit text length
        
        # Format as monospace for better readability
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Courier New'
                run.font.size = Pt(10)


def create_presentation(notebook_path=None, output_dir=None, title=None):
    """
    Convenience function to create a presentation in one call.
    
    Args:
        notebook_path: Path to notebook (auto-detected if None)
        output_dir: Output directory (default: output/presentations/)
        title: Presentation title (auto-generated if None)
        
    Returns:
        Path to generated presentation
    """
    presenter = NotebookPresenter(notebook_path, output_dir)
    presenter.auto_extract_from_notebook()
    return presenter.generate_presentation(title)


if __name__ == '__main__':
    # Example usage
    output_path = create_presentation()
    print(f"Presentation created: {output_path}")
